#!/usr/bin/env python3
"""
fill_template.py — NBB Report PPTX Filler
Remplit le template T21_HK_Agencies_Glass avec les données Excel.
"""

import sys, os, io
import pandas as pd
from pptx import Presentation

AGENCY_GROUP = {
    'SPARK FOUNDRY':    'Publicis Media',
    'STARCOM':          'Publicis Media',
    'ZENITH':           'Publicis Media',
    'PUBLICIS MEDIA':   'Publicis Media',
    'PHD':              'Omnicom Media',
    'OMD':              'Omnicom Media',
    'UM':               'Omnicom Media',
    'INITIATIVE':       'Omnicom Media',
    'HEARTS & SCIENCE': 'Omnicom Media',
    'CARAT':            'Dentsu',
    'IPROSPECT':        'Dentsu',
    'DENTSU X':         'Dentsu',
    'HAVAS MEDIA':      'Havas Media Network',
    'ARENA':            'Havas Media Network',
    'ESSENCEMEDIACOM':  'WPP Media',
    'MINDSHARE':        'WPP Media',
    'WAVEMAKER':        'WPP Media',
    'VALE MEDIA':       'Independant',
}

GROUP_ORDER = [
    'Publicis Media', 'Omnicom Media', 'Dentsu',
    'Havas Media Network', 'WPP Media', 'Independant',
]

GROUP_KEYS = {
    'Publicis Media':       'PUBLICIS',
    'Omnicom Media':        'OMNICOM',
    'Dentsu':               'DENTSU',
    'Havas Media Network':  'HAVAS',
    'WPP Media':            'WPP',
    'Independant':          'INDEP',
}

MAX_AGENCIES   = 14
MAX_WINS_DEPS  = 5
MAX_RETS_S3    = 4
MAX_RETS_S6    = 8
CONCAT_MAX     = 3
CONCAT_MIN_VAL = 0

def fmt_nbb(v):
    if v > 0:   return f'+{v:.1f}m$'
    elif v < 0: return f'{v:.1f}m$'
    return '$0m'

def fmt_val(v):
    if v > 0:   return f'+{v:.1f}m'
    elif v < 0: return f'{v:.1f}m'
    return ''

def trunc(text, limit=25):
    text = str(text).strip()
    return text[:limit-1] + '…' if len(text) > limit else text

def concat_moves(rows, min_val=CONCAT_MIN_VAL, max_items=CONCAT_MAX):
    items = []
    for r in rows:
        val = float(r.get('Integrated Spends', 0) or 0)
        if abs(val) >= min_val and r.get('Advertiser'):
            adv = trunc(str(r['Advertiser']), 20)
            v   = fmt_val(val)
            items.append(f'{adv} {v}'.strip() if v else adv)
        if len(items) >= max_items:
            break
    return '  ·  '.join(items)

def load_data(excel_path):
    df = pd.read_excel(excel_path)
    return load_data_from_df(df)

def load_data_from_df(df):
    df = df.copy()
    df['Agency']            = df['Agency'].astype(str).str.strip().str.upper()
    df['NewBiz']            = df['NewBiz'].astype(str).str.strip().str.upper()
    df['Integrated Spends'] = pd.to_numeric(df['Integrated Spends'], errors='coerce').fillna(0)

    W_all = df[df.NewBiz == 'WIN']
    D_all = df[df.NewBiz == 'DEPARTURE']
    R_all = df[df.NewBiz == 'RETENTION']

    agencies = []
    for ag in df['Agency'].dropna().unique():
        if not ag or ag == 'NAN':
            continue
        sub = df[df['Agency'] == ag]
        W = sub[sub.NewBiz == 'WIN']
        D = sub[sub.NewBiz == 'DEPARTURE']
        R = sub[sub.NewBiz == 'RETENTION']
        w = float(W['Integrated Spends'].sum())
        d = float(D['Integrated Spends'].sum())
        r = float(R['Integrated Spends'].sum())
        agencies.append({
            'agency':    ag,
            'group':     AGENCY_GROUP.get(ag, 'Independant'),
            'nbb':       w + d,
            'wins':      w,
            'deps':      d,
            'rets':      r,
            'wc':        len(W),
            'dc':        len(D),
            'rc':        len(R),
            'wins_rows': W.sort_values('Integrated Spends', ascending=False).to_dict('records'),
            'dep_rows':  D.sort_values('Integrated Spends').to_dict('records'),
            'ret_rows':  R.sort_values('Integrated Spends', ascending=False).to_dict('records'),
        })
    agencies.sort(key=lambda x: -x['nbb'])
    for i, a in enumerate(agencies):
        a['rank'] = i + 1

    # Groupes — rang calculé dynamiquement par NBB
    group_stats = {}
    for g in GROUP_ORDER:
        ags = [a for a in agencies if a['group'] == g]
        group_stats[g] = {
            'name':     g,
            'nbb':      sum(a['nbb']  for a in ags),
            'wins':     sum(a['wins'] for a in ags),
            'deps':     sum(a['deps'] for a in ags),
            'rets':     sum(a['rets'] for a in ags),
            'wc':       sum(a['wc']   for a in ags),
            'dc':       sum(a['dc']   for a in ags),
            'agencies': ags,
        }
    for rank_i, gs in enumerate(
            sorted(group_stats.values(), key=lambda x: -x['nbb']), 1):
        gs['rank'] = rank_i

    top_wins = W_all.sort_values('Integrated Spends', ascending=False).head(MAX_WINS_DEPS).to_dict('records')
    top_deps = D_all.sort_values('Integrated Spends').head(MAX_WINS_DEPS).to_dict('records')
    top_rets = R_all.sort_values('Integrated Spends', ascending=False).head(MAX_RETS_S3).to_dict('records')

    ret_by_agency = []
    for a in agencies:
        if a['rets'] != 0 or a['rc'] > 0:
            top_ret_client = a['ret_rows'][0]['Advertiser'] if a['ret_rows'] else '—'
            ret_by_agency.append({
                'agency':     a['agency'],
                'balance':    a['rets'],
                'top_client': trunc(str(top_ret_client), 25),
            })
    ret_by_agency.sort(key=lambda x: -x['balance'])

    return {
        'agencies':      agencies,
        'group_stats':   group_stats,
        'top_wins':      top_wins,
        'top_deps':      top_deps,
        'top_rets':      top_rets,
        'ret_by_agency': ret_by_agency,
    }

def build_placeholders(data):
    ph = {}
    agencies    = data['agencies']
    group_stats = data['group_stats']
    top_wins    = data['top_wins']
    top_deps    = data['top_deps']
    top_rets    = data['top_rets']
    ret_by_ag   = data['ret_by_agency']

    # Slide 03
    for i in range(1, MAX_WINS_DEPS + 1):
        idx = i - 1
        if idx < len(top_wins):
            r = top_wins[idx]
            ph[f'{{{{WIN_ADV_{i}}}}}'] = trunc(str(r.get('Advertiser', '')), 28)
            ph[f'{{{{WIN_AG_{i}}}}}']  = trunc(str(r.get('Agency', '')), 22)
            ph[f'{{{{WIN_VAL_{i}}}}}'] = fmt_val(float(r.get('Integrated Spends', 0)))
        else:
            ph[f'{{{{WIN_ADV_{i}}}}}'] = ph[f'{{{{WIN_AG_{i}}}}}'] = ph[f'{{{{WIN_VAL_{i}}}}}'] = ''
        if idx < len(top_deps):
            r = top_deps[idx]
            ph[f'{{{{DEP_ADV_{i}}}}}'] = trunc(str(r.get('Advertiser', '')), 28)
            ph[f'{{{{DEP_AG_{i}}}}}']  = trunc(str(r.get('Agency', '')), 22)
            ph[f'{{{{DEP_VAL_{i}}}}}'] = fmt_val(float(r.get('Integrated Spends', 0)))
        else:
            ph[f'{{{{DEP_ADV_{i}}}}}'] = ph[f'{{{{DEP_AG_{i}}}}}'] = ph[f'{{{{DEP_VAL_{i}}}}}'] = ''

    for i in range(1, MAX_RETS_S3 + 1):
        idx = i - 1
        if idx < len(top_rets):
            r = top_rets[idx]
            ph[f'{{{{RET_{i}}}}}']    = trunc(str(r.get('Advertiser', '')), 25)
            ph[f'{{{{RET_AG_{i}}}}}'] = trunc(str(r.get('Agency', '')), 22)
        else:
            ph[f'{{{{RET_{i}}}}}'] = ph[f'{{{{RET_AG_{i}}}}}'] = ''

    # Slide 02
    for i in range(1, 5):
        idx = i - 1
        if idx < len(agencies):
            a = agencies[idx]
            ph[f'{{{{AG_{i}}}}}']       = trunc(a['agency'].title(), 22)
            ph[f'{{{{NBB_{i}}}}}']      = fmt_nbb(a['nbb'])
            ph[f'{{{{WINS_RAW_{i}}}}}'] = str(int(round(a['wins'])))
            ph[f'{{{{DEPS_RAW_{i}}}}}'] = str(int(round(a['deps'])))
        else:
            ph[f'{{{{AG_{i}}}}}'] = ph[f'{{{{NBB_{i}}}}}'] = ''
            ph[f'{{{{WINS_RAW_{i}}}}}'] = ph[f'{{{{DEPS_RAW_{i}}}}}'] = ''

    top3 = agencies[:3]
    lines = []
    for a in top3:
        lines.append(f'• {a["agency"].title()} : NBB {fmt_nbb(a["nbb"])}  (W={fmt_val(a["wins"])} / D={fmt_val(a["deps"])})')
    ph['{{KEY_TAKEAWAYS}}'] = '\n'.join(lines)

    # Slide 04
    for i in range(1, MAX_AGENCIES + 1):
        idx = i - 1
        if idx < len(agencies):
            a = agencies[idx]
            ph[f'{{{{AG_{i}}}}}']      = trunc(a['agency'], 22)
            ph[f'{{{{NBB_{i}}}}}']     = fmt_nbb(a['nbb'])
            ph[f'{{{{RANK_{i}}}}}']    = f'#{a["rank"]}'
            ph[f'{{{{WIN_{i}}}}}']     = f'{a["wins"]:.1f}' if a['wins'] else '0'
            ph[f'{{{{DEP_{i}}}}}']     = f'{a["deps"]:.1f}' if a['deps'] else '0'
            ph[f'{{{{TOPWINS_{i}}}}}'] = concat_moves(a['wins_rows'])
            ph[f'{{{{TOPDEPS_{i}}}}}'] = concat_moves(a['dep_rows'])
            top_ret = trunc(str(a['ret_rows'][0]['Advertiser']), 22) if a['ret_rows'] else '—'
            ph[f'{{{{TOPRET_{i}}}}}']  = top_ret
        else:
            for k in ['AG', 'NBB', 'RANK', 'WIN', 'DEP', 'TOPWINS', 'TOPDEPS', 'TOPRET']:
                ph[f'{{{{{k}_{i}}}}}'] = ''

    # Slide 05 — groupes, rang dynamique
    for g, gs in group_stats.items():
        key = GROUP_KEYS.get(g, g.upper().replace(' ', '_'))
        ph[f'{{{{GRP_RANK_{key}}}}}']  = f'#{gs["rank"]}'
        ph[f'{{{{GRP_NAME_{key}}}}}']  = gs['name']
        ph[f'{{{{GRP_NBB_{key}}}}}']   = fmt_nbb(gs['nbb'])
        ph[f'{{{{GRP_WINS_{key}}}}}']  = f'{gs["wins"]:.1f}'
        ph[f'{{{{GRP_DEPS_{key}}}}}']  = f'{gs["deps"]:.1f}'
        ph[f'{{{{GRP_WC_{key}}}}}']    = str(gs['wc'])
        ph[f'{{{{GRP_DC_{key}}}}}']    = str(gs['dc'])
        MAX_PER_GROUP = 5  # max agencies per group in the template
        for j in range(1, MAX_PER_GROUP + 1):
            idx = j - 1
            if idx < len(gs['agencies']):
                a = gs['agencies'][idx]
                ph[f'{{{{GRP_AG_{key}_{j}}}}}']  = trunc(a['agency'], 20)
                ph[f'{{{{GRP_NBB_{key}_{j}}}}}'] = fmt_nbb(a['nbb'])
                ph[f'{{{{GRP_WIN_{key}_{j}}}}}'] = f'{a["wins"]:.1f}'
                ph[f'{{{{GRP_DEP_{key}_{j}}}}}'] = f'{a["deps"]:.1f}'
                ph[f'{{{{GRP_WC_{key}_{j}}}}}']  = str(a['wc'])
                ph[f'{{{{GRP_DC_{key}_{j}}}}}']  = str(a['dc'])
            else:
                # Clear unused agency rows in the template
                for k in ['GRP_AG', 'GRP_NBB', 'GRP_WIN', 'GRP_DEP', 'GRP_WC', 'GRP_DC']:
                    ph[f'{{{{{k}_{key}_{j}}}}}'] = ''

    # Slide 06
    for i in range(1, MAX_RETS_S6 + 1):
        idx = i - 1
        if idx < len(ret_by_ag):
            r = ret_by_ag[idx]
            ph[f'{{{{RET_AG_{i}}}}}']  = trunc(r['agency'], 22)
            ph[f'{{{{RET_BAL_{i}}}}}'] = fmt_nbb(r['balance'])
            ph[f'{{{{RET_TOP_{i}}}}}'] = r['top_client']
        else:
            ph[f'{{{{RET_AG_{i}}}}}'] = ph[f'{{{{RET_BAL_{i}}}}}'] = ph[f'{{{{RET_TOP_{i}}}}}'] = ''

    return ph

def _replace_in_paragraph(para, replacements):
    full_text = ''.join(run.text for run in para.runs)
    if '{{' not in full_text:
        return
    new_text = full_text
    for placeholder, value in replacements.items():
        new_text = new_text.replace(placeholder, str(value))
    if new_text == full_text:
        return
    if para.runs:
        para.runs[0].text = new_text
        for run in para.runs[1:]:
            run.text = ''

def replace_all_placeholders(prs, replacements):
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    _replace_in_paragraph(para, replacements)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            _replace_in_paragraph(para, replacements)

def fill_template(excel_path, template_path, output_path):
    print(f'📂 Excel    : {excel_path}')
    print(f'📐 Template : {template_path}')
    data = load_data(excel_path)
    print(f'✅ {len(data["agencies"])} agences chargées')
    ph = build_placeholders(data)
    print(f'✅ {len(ph)} balises construites')
    prs = Presentation(template_path)
    replace_all_placeholders(prs, ph)
    print('✅ Balises remplacées')
    os.makedirs(os.path.dirname(output_path) or '.', exist_ok=True)
    prs.save(output_path)
    print(f'✅ Sauvegardé : {output_path}  ({os.path.getsize(output_path)//1024} KB)')
    return output_path

if __name__ == '__main__':
    EXCEL    = sys.argv[1] if len(sys.argv) > 1 else '/mnt/user-data/uploads/Newbiz_Balance_DB_Report__1_.xlsx'
    TEMPLATE = sys.argv[2] if len(sys.argv) > 2 else 'T21_HK_Agencies_Glass_v13.pptx'
    OUTPUT   = sys.argv[3] if len(sys.argv) > 3 else '/mnt/user-data/outputs/NBB_filled.pptx'
    fill_template(EXCEL, TEMPLATE, OUTPUT)
