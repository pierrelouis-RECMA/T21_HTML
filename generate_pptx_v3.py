#!/usr/bin/env python3
"""
generate_pptx_v3.py — NBB Report Generator (unified pipeline)
───────────────────────────────────────────────────────────
Lit un DataFrame pandas (issu de l'Excel NBB) et génère
les slides "Agency Cards" en XML direct dans le template
T21_HK_Agencies_Glass_v12.pptx

API publique (appelée depuis app.py) :
    build_agency_pptx(df, template_path) -> bytes  (PPTX en mémoire)

Usage standalone (test local) :
    python3 generate_pptx_v2.py <excel> <template> <output>

════════════════════════════════════════════════════════════
COLONNES EXCEL ATTENDUES
════════════════════════════════════════════════════════════

OBLIGATOIRES — la génération échoue si ces colonnes sont absentes :

  Agency               str   Nom de l'agence, ex: "MINDSHARE", "HAVAS MEDIA"
                             → Mis en majuscules automatiquement.
                             → Utilisé pour le titre de la carte et le
                               classement NBB.

  NewBiz               str   Type de mouvement. Valeurs acceptées (insensible
                             à la casse) :
                               WIN       = l'agence gagne ce client
                               DEPARTURE = l'agence perd ce client
                               RETENTION = renouvellement / rétention
                             → Toute autre valeur est ignorée.

  Advertiser           str   Nom de l'annonceur, ex: "RECKITT", "MARS WRIGLEY"
                             → Affiché dans la colonne WIN/DEPARTURE/RETENTION
                               de la carte. Tronqué à 28 caractères.

  Integrated Spends    float Budget en millions USD.
                             → Positif pour les WIN (ex: 150.0)
                             → Négatif pour les DEPARTURE (ex: -28.3)
                             → Utilisé pour calculer le NBB de chaque agence :
                               NBB = somme(WIN) + somme(DEPARTURE)
                             → Les RETENTION ne comptent PAS dans le NBB.
                             → Affiché comme "+150.0m" ou "-28.3m" sur la carte.

OPTIONNELLES — ignorées si absentes, utilisées si présentes :

  Date of announcement date  Date d'annonce du mouvement.
                             → Affiché au format "Jan-25" sur la carte.
                             → Accepte timestamp Excel, string ISO, ou datetime.

  Incumbent            str   Agence précédente (pour les WIN uniquement).
                             → Non affiché dans les cards actuellement,
                               disponible pour extensions futures.

IGNORÉES (présentes dans l'Excel NBB standard, non utilisées) :
  Status, Country, Brand, "% that you handle", Years,
  "Date of Effectiveness", "Country of Decision", "Move ?",
  "Pitch participation ?", Assignment, Contender, Remarks,
  "Share of Non trad", "Ad Spends", "Ask Agency for validation"

════════════════════════════════════════════════════════════
LOGIQUE DE GÉNÉRATION
════════════════════════════════════════════════════════════

1. Les agences sont triées par NBB décroissant.
2. Elles sont réparties en slides de AGENCIES_PER_SLIDE (défaut: 4) agences.
   → 16 agences = 4 slides (22, 23, 24, 25)
   → 20 agences = 5 slides (22 … 26)
   → Aucune limite — s'adapte automatiquement.
3. Le numéro de slide démarre à 22 (slides 1-21 = autres sections du deck).
4. Les agences inconnues dans AGENCY_GROUP sont classées "Independant".

════════════════════════════════════════════════════════════
AJOUTER UN NOUVEAU MARCHÉ (ex: Singapore)
════════════════════════════════════════════════════════════

1. Ajouter les agences Singapore dans AGENCY_GROUP ci-dessous.
2. Uploader l'Excel Singapore via l'interface web.
3. Si les noms de colonnes diffèrent, utiliser le mapping de colonnes
   dans l'interface (affiché automatiquement).
4. Aucun autre changement nécessaire — tout est dynamique.
"""

import io, os, re, shutil, subprocess, tempfile, zipfile
import xml.etree.ElementTree as ET
import pandas as pd

# ── XML namespaces ────────────────────────────────────────────
A  = "http://schemas.openxmlformats.org/drawingml/2006/main"
P  = "http://schemas.openxmlformats.org/presentationml/2006/main"
ET.register_namespace('a', A)
ET.register_namespace('p', P)
ET.register_namespace('r', "http://schemas.openxmlformats.org/officeDocument/2006/relationships")
ta = lambda l: f'{{{A}}}{l}'
tp = lambda l: f'{{{P}}}{l}'

# ── Slide geometry (A4 portrait) ─────────────────────────────
SW, SH    = 7559675, 10439400
NAV_H     = 370000
MARG_X    = 220000
MARG_TOP  = 100000
FOOTER_Y  = SH - 170000 - 30000
CARD_W    = SW - 2 * MARG_X
CARD_GAP  = 100000

# ── Design tokens ─────────────────────────────────────────────
C_HEADER_BG  = "0F172A"
C_HEADER_TXT = "FFFFFF"
C_WIN_TXT    = "059669"
C_WIN_UL     = "10B981"
C_DEP_TXT    = "BE123C"
C_DEP_UL     = "F43F5E"
C_RET_TXT    = "B45309"
C_RET_UL     = "F59E0B"
C_CARD_BG    = "FFFFFF"
C_COL_BG     = "F8FAFC"
C_ITEM_TXT   = "334155"
C_VAL_POS    = "059669"
C_VAL_NEG    = "E11D48"
C_EMPTY      = "94A3B8"
C_BORDER     = "E2E8F0"
C_NAV_BG     = "0F172A"
C_NAV_TXT    = "64748B"
C_FOOT       = "94A3B8"
FONT         = "Segoe UI"

NAV_LABELS   = ["Key Findings", "TOP moves", "NBB · Agencies",
                "NBB · Groups", "Retentions", "Details"]

# ── Group / market config ─────────────────────────────────────
AGENCY_GROUP = {
    # Publicis Media
    "SPARK FOUNDRY":    "Publicis Media",
    "STARCOM":          "Publicis Media",
    "ZENITH":           "Publicis Media",
    "PUBLICIS MEDIA":   "Publicis Media",
    # Omnicom Media Group
    "PHD":              "Omnicom Media",
    "OMD":              "Omnicom Media",
    "UM":               "Omnicom Media",
    "INITIATIVE":       "Omnicom Media",
    "HEARTS & SCIENCE": "Omnicom Media",
    # Dentsu
    "CARAT":            "Dentsu",
    "IPROSPECT":        "Dentsu",
    "DENTSU X":         "Dentsu",
    # Havas Media Network
    "HAVAS MEDIA":      "Havas Media Network",
    "ARENA":            "Havas Media Network",
    # WPP Media
    "ESSENCEMEDIACOM":  "WPP Media",
    "MINDSHARE":        "WPP Media",
    "WAVEMAKER":        "WPP Media",
    # Independants
    "VALE MEDIA":       "Independant",
}

AGENCIES_PER_SLIDE = 4   # maximum par slide — ajoute des slides si besoin

# ── ID counter (reset par slide) ─────────────────────────────
_id = [1000]
def nid():
    _id[0] += 1
    return _id[0]

# ─────────────────────────────────────────────────────────────
# 1. DATA LAYER — Excel → structured dicts
# ─────────────────────────────────────────────────────────────

def fmt_date(v):
    try:
        if isinstance(v, (int, float)):
            d = pd.to_datetime(v, origin="1899-12-30", unit="D")
        else:
            d = pd.to_datetime(v)
        return d.strftime("%b-%y")
    except:
        return ""

def fmt_nbb(v: float) -> str:
    if v > 0:   return f"+{v:.1f}m$"
    elif v < 0: return f"{v:.1f}m$"
    else:       return "$0m"

def make_items(rows_df: pd.DataFrame) -> list:
    items = []
    for _, r in rows_df.iterrows():
        adv = str(r.get("Advertiser", "")).strip()[:28]
        val = float(r.get("Integrated Spends", 0) or 0)
        val_str = f"+{val:.1f}m" if val > 0 else (f"{val:.1f}m" if val != 0 else "")
        items.append((adv, val_str))
    return items

def df_to_agencies(df: pd.DataFrame) -> list:
    """
    Prend un DataFrame NBB (colonnes Agency, NewBiz, Advertiser,
    Integrated Spends, Date of announcement, Incumbent)
    et retourne une liste de dicts agency triée par NBB desc.
    """
    df = df.copy()
    df["Agency"] = df["Agency"].astype(str).str.strip().str.upper()
    df["NewBiz"]  = df["NewBiz"].astype(str).str.strip().str.upper()
    df["Integrated Spends"] = pd.to_numeric(df["Integrated Spends"], errors="coerce").fillna(0)

    result = []
    for ag in df["Agency"].dropna().unique():
        if not ag or ag == "NAN":
            continue
        sub = df[df["Agency"] == ag]
        W = sub[sub.NewBiz == "WIN"]
        D = sub[sub.NewBiz == "DEPARTURE"]
        R = sub[sub.NewBiz == "RETENTION"]
        w_sum = float(W["Integrated Spends"].sum())
        d_sum = float(D["Integrated Spends"].sum())
        nbb   = w_sum + d_sum

        result.append({
            "name":     ag,
            "group":    AGENCY_GROUP.get(ag, "Independant"),
            "nbb":      fmt_nbb(nbb),
            "_nbb_raw": nbb,
            "wins":     make_items(W),
            "deps":     make_items(D),
            "rets":     make_items(R),
        })

    result.sort(key=lambda x: -x["_nbb_raw"])
    return result

def agencies_to_slides(agencies: list) -> dict:
    """
    Distribue les agences en slides de AGENCIES_PER_SLIDE max.
    Retourne {slide_num: [agency_dict, ...]}
    """
    slides = {}
    for i, start in enumerate(range(0, len(agencies), AGENCIES_PER_SLIDE)):
        slides[7 + i] = agencies[start:start + AGENCIES_PER_SLIDE]
    return slides

# ─────────────────────────────────────────────────────────────
# 2. XML BUILDERS
# ─────────────────────────────────────────────────────────────

def make_shadow():
    eff  = ET.Element(ta("effectLst"))
    shdw = ET.SubElement(eff, ta("outerShdw"))
    shdw.set("blurRad", "38100")
    shdw.set("dist",    "12700")
    shdw.set("dir",     "5400000")
    shdw.set("algn",    "ctr")
    shdw.set("rotWithShape", "0")
    clr = ET.SubElement(shdw, ta("srgbClr"))
    clr.set("val", "1A2B4A")
    ET.SubElement(clr, ta("alpha")).set("val", "10000")
    return eff

def mk_sp(id_, name, x, y, w, h, geom="rect", fill=None, no_fill=False,
          corner=None, border_c=None, border_w=4762, shadow=False):
    sp = ET.Element(tp("sp"))
    nv = ET.SubElement(sp, tp("nvSpPr"))
    cn = ET.SubElement(nv, tp("cNvPr"))
    cn.set("id", str(id_)); cn.set("name", name)
    ET.SubElement(nv, tp("cNvSpPr")); ET.SubElement(nv, tp("nvPr"))
    pr = ET.SubElement(sp, tp("spPr"))
    xf = ET.SubElement(pr, ta("xfrm"))
    o  = ET.SubElement(xf, ta("off")); o.set("x", str(x)); o.set("y", str(y))
    e  = ET.SubElement(xf, ta("ext")); e.set("cx", str(w)); e.set("cy", str(h))
    pg = ET.SubElement(pr, ta("prstGeom")); pg.set("prst", geom)
    av = ET.SubElement(pg, ta("avLst"))
    if geom == "roundRect" and corner is not None:
        gd = ET.SubElement(av, ta("gd"))
        gd.set("name", "adj"); gd.set("fmla", f"val {corner}")
    if no_fill:
        ET.SubElement(pr, ta("noFill"))
    elif fill:
        sf = ET.SubElement(pr, ta("solidFill"))
        ET.SubElement(sf, ta("srgbClr")).set("val", fill)
    ln = ET.SubElement(pr, ta("ln"))
    if border_c:
        ln.set("w", str(border_w))
        sf2 = ET.SubElement(ln, ta("solidFill"))
        ET.SubElement(sf2, ta("srgbClr")).set("val", border_c)
    else:
        ET.SubElement(ln, ta("noFill"))
    if shadow:
        pr.append(make_shadow())
    tb = ET.SubElement(sp, tp("txBody"))
    ET.SubElement(tb, ta("bodyPr"))
    ET.SubElement(tb, ta("lstStyle"))
    ET.SubElement(ET.SubElement(tb, ta("p")), ta("endParaRPr")).set("lang", "en-US")
    return sp

def mk_tx(id_, name, x, y, w, h, runs_or_paras,
          anchor="t", algn="l", lIns=0, rIns=0, tIns=0, bIns=0,
          autofit=False, spc=0):
    if runs_or_paras and not isinstance(runs_or_paras[0], (list, tuple)):
        paras = [runs_or_paras]
    elif runs_or_paras and isinstance(runs_or_paras[0][0], str):
        paras = [runs_or_paras]
    else:
        paras = runs_or_paras

    sp = ET.Element(tp("sp"))
    nv = ET.SubElement(sp, tp("nvSpPr"))
    cn = ET.SubElement(nv, tp("cNvPr")); cn.set("id", str(id_)); cn.set("name", name)
    cs = ET.SubElement(nv, tp("cNvSpPr")); cs.set("txBox", "1")
    ET.SubElement(nv, tp("nvPr"))
    pr = ET.SubElement(sp, tp("spPr"))
    xf = ET.SubElement(pr, ta("xfrm"))
    o  = ET.SubElement(xf, ta("off")); o.set("x", str(x)); o.set("y", str(y))
    e  = ET.SubElement(xf, ta("ext")); e.set("cx", str(w)); e.set("cy", str(h))
    ET.SubElement(ET.SubElement(pr, ta("prstGeom")), ta("avLst"))
    pr.find(ta("prstGeom")).set("prst", "rect")
    ET.SubElement(pr, ta("noFill"))
    ET.SubElement(ET.SubElement(pr, ta("ln")), ta("noFill"))
    tb = ET.SubElement(sp, tp("txBody"))
    bp = ET.SubElement(tb, ta("bodyPr"))
    bp.set("wrap", "square"); bp.set("anchor", anchor)
    bp.set("lIns", str(lIns)); bp.set("rIns", str(rIns))
    bp.set("tIns", str(tIns)); bp.set("bIns", str(bIns))
    if autofit:
        ET.SubElement(bp, ta("normAutofit"))
    ET.SubElement(tb, ta("lstStyle"))

    for para in paras:
        p_el = ET.SubElement(tb, ta("p"))
        pp   = ET.SubElement(p_el, ta("pPr")); pp.set("algn", algn)
        if not para:
            ET.SubElement(p_el, ta("endParaRPr")).set("lang", "en-US")
            continue
        for run in para:
            txt   = run[0]
            sz    = run[1] if len(run) > 1 else 900
            bold  = run[2] if len(run) > 2 else False
            color = run[3] if len(run) > 3 else C_ITEM_TXT
            ital  = run[4] if len(run) > 4 else False
            r_spc = run[5] if len(run) > 5 else spc
            r_el  = ET.SubElement(p_el, ta("r"))
            rp    = ET.SubElement(r_el, ta("rPr"))
            rp.set("lang", "en-US"); rp.set("sz", str(sz))
            rp.set("b", "1" if bold else "0")
            if ital: rp.set("i", "1")
            if r_spc: rp.set("spc", str(r_spc))
            rp.set("dirty", "0")
            sf = ET.SubElement(rp, ta("solidFill"))
            ET.SubElement(sf, ta("srgbClr")).set("val", color)
            ET.SubElement(rp, ta("latin")).set("typeface", FONT)
            ET.SubElement(r_el, ta("t")).text = txt
    return sp

# ── Nav bar ───────────────────────────────────────────────────
def nav_bar(active=5, nav_rids=None):
    """
    nav_rids : dict {section_index: rId_string} pour les liens hypertexte.
               section_index = 0..5 (Key Findings ... Details)
               Si None, les rectangles invisibles ne sont pas generés.
    """
    out = [mk_sp(nid(), "NavBg", 0, 0, SW, NAV_H, fill=C_NAV_BG)]
    cw  = SW // 6
    for i, lbl in enumerate(NAV_LABELS):
        act = i == active
        if i:
            out.append(mk_sp(nid(), f"NS{i}", i*cw - 3000, 35000, 6000, 280000, fill="1E293B"))
        if act:
            out.append(mk_sp(nid(), f"NA{i}", i*cw, 0, cw, NAV_H, fill="1E293B"))
        out.append(mk_tx(nid(), f"NN{i}", i*cw, 22000, cw, 120000,
                         [("0"+str(i+1), 750, False, "FFFFFF" if act else "475569", False, 50)],
                         algn="c"))
        out.append(mk_tx(nid(), f"NL{i}", i*cw, 155000, cw, 165000,
                         [(lbl, 800, act, "FFFFFF" if act else C_NAV_TXT, False, 20)],
                         algn="c"))
    out.append(mk_sp(nid(), "NAcc", 0, NAV_H - 3000, SW, 3000, fill="1E293B"))

    # Rectangles invisibles avec liens (reproduit exactement le template)
    # Positions mesurees sur les slides agency du template (slide7-10)
    # La section active ne recoit pas de lien (pas de lien vers soi-meme)
    NAV_HITBOXES = [
        (0,   92583,  30000, 1152362, 310000),  # Key Findings -> slide2
        (1, 1526883,  12499,  881069, 325000),  # TOP moves    -> slide3
        (2, 2699891,  29583, 1152362, 310000),  # NBB Agencies -> slide4
        (3, 4037822,  20270, 1152362, 310000),  # NBB Groups   -> slide5
        (4, 5284791,  29583,  966597, 310000),  # Retentions   -> slide6
        (5, 6572720,  29583,  909821, 310000),  # Details      -> 1re slide agency
    ]
    if nav_rids:
        R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
        for sec_idx, x, y, w, h in NAV_HITBOXES:
            # On skip uniquement si pas de rId (= make_nav_rid_map a exclu cette cible
            # car elle pointe vers la slide courante elle-meme)
            rid = nav_rids.get(sec_idx)
            if not rid:
                continue
            sp_id = nid()
            sp = ET.Element(tp("sp"))
            nv = ET.SubElement(sp, tp("nvSpPr"))
            cn = ET.SubElement(nv, tp("cNvPr"))
            cn.set("id", str(sp_id))
            cn.set("name", f"NavLink{sec_idx}")
            hl = ET.SubElement(cn, ta("hlinkClick"))
            hl.set(f"{{{R_NS}}}id", rid)
            hl.set("action", "ppaction://hlinksldjump")
            ET.SubElement(nv, tp("cNvSpPr"))
            ET.SubElement(nv, tp("nvPr"))
            pr = ET.SubElement(sp, tp("spPr"))
            xf = ET.SubElement(pr, ta("xfrm"))
            o  = ET.SubElement(xf, ta("off")); o.set("x", str(x)); o.set("y", str(y))
            e  = ET.SubElement(xf, ta("ext")); e.set("cx", str(w)); e.set("cy", str(h))
            pg = ET.SubElement(pr, ta("prstGeom")); pg.set("prst", "roundRect")
            ET.SubElement(pg, ta("avLst"))
            ET.SubElement(pr, ta("noFill"))
            ln = ET.SubElement(pr, ta("ln"))
            ET.SubElement(ln, ta("noFill"))
            tb = ET.SubElement(sp, tp("txBody"))
            ET.SubElement(tb, ta("bodyPr"))
            ET.SubElement(tb, ta("lstStyle"))
            ET.SubElement(ET.SubElement(tb, ta("p")), ta("endParaRPr")).set("lang", "en-US")
            out.append(sp)

    return out

def footer():
    return [mk_tx(nid(), "Footer", MARG_X, FOOTER_Y, CARD_W, 160000,
                  [("Détail par agence  ·  Retentions & renewals not included in NBB calculation",
                    750, False, C_FOOT, False, 20)],
                  algn="c")]

# ── Card height calculator ────────────────────────────────────
PAD_H          = 220000
PAD_V          = 160000
HDR_H          = 290000
SEP_OFFSET     = PAD_V + HDR_H + 60000
BODY_OFFSET    = SEP_OFFSET + 9000 + 120000
COL_HDR_H      = 160000
UL_H           = 20000
ITEM_START_OFF = BODY_OFFSET + 130000 + COL_HDR_H + 20000 + UL_H + 100000
BOT_PAD        = PAD_V
FIXED_OVERHEAD = ITEM_START_OFF + BOT_PAD

def card_heights(agencies: list, available: int) -> list:
    n     = len(agencies)
    space = available - (n - 1) * CARD_GAP

    def max_ni(ag):
        return max(len(ag["wins"]), len(ag["deps"]), len(ag["rets"]), 1)

    if n == 1:
        ni  = max_ni(agencies[0])
        h   = FIXED_OVERHEAD + ni * 190000
        cap = int(space * 0.35)
        return [min(h, cap)]

    n_items = [max_ni(ag) for ag in agencies]
    bases   = [FIXED_OVERHEAD + ni * 155000 for ni in n_items]
    total   = sum(bases)
    hs      = [int(b * space / total) for b in bases]
    hs[-1] += space - sum(hs)
    return hs

# ── Agency card ───────────────────────────────────────────────
def agency_card(ag: dict, bx, by, bw, bh, pfx: str) -> list:
    out = []
    out.append(mk_sp(nid(), f"{pfx}Bg", bx, by, bw, bh,
                     geom="roundRect", fill=C_CARD_BG,
                     corner=4000, border_c=C_BORDER, border_w=4762, shadow=True))

    # Header
    nm_x    = bx + PAD_H
    badge_h = 220000
    badge_w = max(1200000, len(ag["nbb"]) * 85000)
    badge_x = bx + bw - PAD_H - badge_w
    badge_y = by + (HDR_H - badge_h) // 2

    out.append(mk_sp(nid(), f"{pfx}Hdr", bx, by, bw, HDR_H,
                     fill=C_HEADER_BG, border_c=None))

    nlen = len(ag["name"])
    nm_sz = 1500 if nlen <= 14 else (1300 if nlen <= 20 else 1100)
    nm_w  = badge_x - nm_x - 80000
    out.append(mk_tx(nid(), f"{pfx}Nm", nm_x, by + 20000, nm_w, HDR_H - 40000,
                     [(ag["name"], nm_sz, True, C_HEADER_TXT, False, 80)],
                     anchor="ctr", spc=80))
    out.append(mk_tx(nid(), f"{pfx}Grp", nm_x, by + HDR_H - 140000, nm_w, 130000,
                     [(f"({ag['group']})", 750, False, "475569", False, 30)],
                     anchor="ctr"))

    # NBB badge
    is_pos     = ag["nbb"].startswith("+")
    is_neg     = ag["nbb"].startswith("-") and ag["nbb"] != "$0m"
    badge_fill = "064E3B" if is_pos else ("881337" if is_neg else "1E293B")
    out.append(mk_sp(nid(), f"{pfx}BadgeBg", badge_x, badge_y, badge_w, badge_h,
                     geom="roundRect", fill=badge_fill, corner=10000))
    badge_c = "6EE7B7" if is_pos else ("FCA5A5" if is_neg else "94A3B8")
    out.append(mk_tx(nid(), f"{pfx}BadgeTx", badge_x, badge_y, badge_w, badge_h,
                     [(ag["nbb"], 1000, True, badge_c, False, 30)],
                     algn="c", anchor="ctr"))

    # Separator
    sep_y = by + HDR_H
    out.append(mk_sp(nid(), f"{pfx}Sep", bx + PAD_H, sep_y,
                     bw - 2*PAD_H, 9000, fill=C_BORDER))

    # 3-column body
    BODY_Y  = sep_y + 9000 + 120000
    BODY_H  = by + bh - PAD_V - BODY_Y
    COL_PAD = 100000
    col_w   = (bw - 2*PAD_H - 2*COL_PAD) // 3

    cols = [
        ("WIN",       C_WIN_TXT, C_WIN_UL, "wins"),
        ("DEPARTURE", C_DEP_TXT, C_DEP_UL, "deps"),
        ("RETENTION", C_RET_TXT, C_RET_UL, "rets"),
    ]
    for ci, (label, ltxt, lul, key) in enumerate(cols):
        cx    = bx + PAD_H + ci * (col_w + COL_PAD)
        items = ag[key]
        CP    = 90000

        out.append(mk_sp(nid(), f"{pfx}C{ci}Bg", cx, BODY_Y, col_w, BODY_H,
                         geom="roundRect", fill=C_COL_BG,
                         corner=3000, border_c=C_BORDER, border_w=4762))

        H4_Y = BODY_Y + 120000
        out.append(mk_tx(nid(), f"{pfx}C{ci}Lbl", cx + CP, H4_Y,
                         col_w - 2*CP, COL_HDR_H,
                         [(label, 850, True, ltxt, False, 100)], spc=100))

        UL_Y = H4_Y + COL_HDR_H + 15000
        out.append(mk_sp(nid(), f"{pfx}C{ci}UL", cx + CP, UL_Y,
                         col_w - 2*CP, UL_H, fill=lul))

        ITEM_Y0    = UL_Y + UL_H + 100000
        ITEM_AVAIL = BODY_Y + BODY_H - ITEM_Y0 - CP

        if not items:
            out.append(mk_tx(nid(), f"{pfx}C{ci}E",
                             cx + CP, ITEM_Y0, col_w - 2*CP, 200000,
                             [("—", 850, False, C_EMPTY, True)]))
        else:
            item_h  = min(200000, max(120000, ITEM_AVAIL // len(items)))
            item_sz = 870 if item_h >= 170000 else 780

            for ji, item in enumerate(items):
                name = item[0]
                val  = item[1] if len(item) > 1 else ""
                iy   = ITEM_Y0 + ji * item_h
                if len(name) > 25: name = name[:24] + "…"

                name_w = col_w - 2*CP - (580000 if val and val != "—" else 0)
                out.append(mk_tx(nid(), f"{pfx}C{ci}I{ji}N",
                                 cx + CP, iy, name_w, item_h,
                                 [(name, item_sz, False, C_ITEM_TXT)], anchor="ctr"))
                if val and val != "—":
                    vc = C_VAL_POS if val.startswith("+") else C_VAL_NEG
                    out.append(mk_tx(nid(), f"{pfx}C{ci}I{ji}V",
                                     cx + col_w - CP - 570000, iy, 570000, item_h,
                                     [(val, item_sz, True, vc)], algn="r", anchor="ctr"))
    return out

# ── Build one slide XML ───────────────────────────────────────
def build_slide_xml(slide_num: int, agencies: list, nav_rids: dict = None) -> str:
    _id[0] = 1000 + (slide_num - 6) * 5000

    content_top = NAV_H + MARG_TOP
    content_bot = FOOTER_Y - 40000
    available   = content_bot - content_top
    card_hs     = card_heights(agencies, available)

    P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
    A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
    sld = ET.Element(tp("sld"))
    cs  = ET.SubElement(sld, tp("cSld"))
    bg  = ET.SubElement(cs, tp("bg"))
    bgP = ET.SubElement(bg, tp("bgPr"))
    sf  = ET.SubElement(bgP, ta("solidFill"))
    ET.SubElement(sf, ta("srgbClr")).set("val", "FFFFFF")
    ET.SubElement(bgP, ta("effectLst"))

    st  = ET.SubElement(cs, tp("spTree"))
    nvG = ET.SubElement(st, tp("nvGrpSpPr"))
    cn  = ET.SubElement(nvG, tp("cNvPr")); cn.set("id", "1"); cn.set("name", "")
    ET.SubElement(nvG, tp("cNvGrpSpPr")); ET.SubElement(nvG, tp("nvPr"))
    gp  = ET.SubElement(st, tp("grpSpPr"))
    xf  = ET.SubElement(gp, ta("xfrm"))
    for tg in ["off", "ext", "chOff", "chExt"]:
        el = ET.SubElement(xf, ta(tg))
        if tg in ("ext", "chExt"): el.set("cx", "0"); el.set("cy", "0")
        else:                      el.set("x",  "0"); el.set("y",  "0")

    # active=5 = "Details" (section des agency cards)
    for s in nav_bar(active=5, nav_rids=nav_rids): st.append(s)

    cy = content_top
    for i, ag in enumerate(agencies):
        for s in agency_card(ag, MARG_X, cy, CARD_W, card_hs[i], f"A{i}"):
            st.append(s)
        cy += card_hs[i] + CARD_GAP

    for s in footer(): st.append(s)

    # PowerPoint requiert p:clrMapOvr dans chaque slide (sinon "cannot read")
    clr = ET.SubElement(sld, tp("clrMapOvr"))
    ET.SubElement(clr, ta("masterClrMapping"))

    # Indent + serialize
    _indent(sld)
    out = io.StringIO()
    out.write('<?xml version="1.0" encoding="utf-8"?>\n')
    ET.ElementTree(sld).write(out, xml_declaration=False, encoding="unicode")
    return out.getvalue()

def _indent(el, lv=0):
    i = "\n" + "  " * lv
    if len(el):
        if not el.text  or not el.text.strip():  el.text  = i + "  "
        if not el.tail  or not el.tail.strip():  el.tail  = i
        for ch in el: _indent(ch, lv + 1)
        if not ch.tail or not ch.tail.strip():   ch.tail  = i
    else:
        if lv and (not el.tail or not el.tail.strip()): el.tail = i

# ─────────────────────────────────────────────────────────────
# 3. PPTX PACKER
# ─────────────────────────────────────────────────────────────

# Les 6 sections de navigation et leur slide cible dans le template
# section_index → filename relatif de la slide cible
NAV_SECTION_TARGETS = {
    0: "slide2.xml",   # Key Findings
    1: "slide3.xml",   # TOP moves
    2: "slide4.xml",   # NBB Agencies
    3: "slide5.xml",   # NBB Groups
    4: "slide6.xml",   # Retentions
    5: None,           # Details → première slide agency (dynamique)
}

def make_slide_rels(nav_rid_map: dict) -> str:
    """
    Génère le fichier .rels d'une slide agency avec :
      - rId1 : slideLayout (obligatoire)
      - rId2..rId7 : liens vers les 6 sections de navigation
    nav_rid_map : {section_index: (rId, target_filename)}
    """
    REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
    SLD_TYPE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide"
    LAY_TYPE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout"

    lines = [
        f'<?xml version="1.0" encoding="utf-8"?>',
        f'<Relationships xmlns="{REL_NS}">',
        f'  <Relationship Id="rId1" Type="{LAY_TYPE}" Target="../slideLayouts/slideLayout1.xml"/>',
    ]
    for sec_idx, (rid, target) in sorted(nav_rid_map.items(), key=lambda x: x[1][0]):
        lines.append(
            f'  <Relationship Id="{rid}" Type="{SLD_TYPE}" Target="{target}"/>'
        )
    lines.append("</Relationships>")
    return "\n".join(lines) + "\n"

def build_agency_pptx(df: pd.DataFrame, template_path: str,
                      prefilled_prs_bytes: bytes = None) -> bytes:
    """
    Point d'entrée unique — appelé par app.py.

    Pipeline en 2 étapes :
      1. fill_template : remplace les balises {{...}} dans les slides 1-6
      2. injection XML : génère les slides agency cards (7+) dynamiquement

    df            : DataFrame Excel (Agency, NewBiz, Advertiser, Integrated Spends…)
    template_path : chemin vers T21_HK_Agencies_Glass_v13.pptx
    prefilled_prs_bytes : optionnel — PPTX déjà pré-rempli (étape 1 déjà faite)
    Retourne : bytes du PPTX final complet.
    """
    # ── Étape 1 : remplir les balises des slides statiques (1-6) ──────────────
    if prefilled_prs_bytes is None:
        from fill_template import load_data_from_df, build_placeholders, replace_all_placeholders
        from pptx import Presentation as _Prs
        _data = load_data_from_df(df)
        _ph   = build_placeholders(_data)
        _prs  = _Prs(template_path)
        replace_all_placeholders(_prs, _ph)
        _buf  = io.BytesIO()
        _prs.save(_buf)
        prefilled_prs_bytes = _buf.getvalue()
        print(f"  ✅ Étape 1 : {len(_ph)} balises remplacées")
    agencies = df_to_agencies(df)
    slides   = agencies_to_slides(agencies)   # {22: [...], 23: [...], ...}

    # La première slide agency = cible du lien "Details" dans la nav
    first_agency_slide_num = min(slides.keys())
    first_agency_filename  = f"slide{first_agency_slide_num}.xml"

    # Source du PPTX : pré-rempli si disponible, sinon template brut
    if prefilled_prs_bytes is not None:
        tpl_bytes = prefilled_prs_bytes
    else:
        with open(template_path, "rb") as f:
            tpl_bytes = f.read()

    with zipfile.ZipFile(io.BytesIO(tpl_bytes), "r") as zin:
        files = {n: zin.read(n) for n in zin.namelist()}

    FIRST_AGENCY_SLIDE = 7

    existing_slides = sorted(
        [k for k in files if re.match(r"ppt/slides/slide(\d+)\.xml$", k)],
        key=lambda k: int(re.search(r"(\d+)", k.split("/")[-1]).group())
    )
    max_existing = max(
        int(re.search(r"(\d+)", k.split("/")[-1]).group())
        for k in existing_slides
    ) if existing_slides else 21

    # Supprimer slides >= FIRST_AGENCY_SLIDE (et leurs rels)
    for snum in range(FIRST_AGENCY_SLIDE, max_existing + 1):
        files.pop(f"ppt/slides/slide{snum}.xml", None)
        files.pop(f"ppt/slides/_rels/slide{snum}.xml.rels", None)

    # ── Construire le mapping nav pour chaque slide agency ───────
    # rId2..rId7 = les 6 sections ; rId1 = slideLayout (réservé)
    # On résout la cible "Details" dynamiquement vers first_agency_filename
    def make_nav_rid_map(current_slide_num: int) -> dict:
        """
        Retourne {section_index: (rId_string, target_filename)}
        pour toutes les sections SAUF celle qui pointe vers soi-même.
        """
        rid_map    = {}   # section_index -> (rId, filename)
        nav_rids   = {}   # section_index -> rId_string  (pour nav_bar())
        rid_counter = 2   # rId1 = slideLayout
        for sec_idx in range(6):
            target = NAV_SECTION_TARGETS[sec_idx]
            if target is None:
                target = first_agency_filename
            # Ne pas créer de lien vers soi-même
            if target == f"slide{current_slide_num}.xml":
                continue
            rid = f"rId{rid_counter}"
            rid_map[sec_idx]  = (rid, target)
            nav_rids[sec_idx] = rid
            rid_counter += 1
        return rid_map, nav_rids

    # Ajouter les nouvelles slides agency
    for snum, ags in slides.items():
        rid_map, nav_rids = make_nav_rid_map(snum)
        xml_str  = build_slide_xml(snum, ags, nav_rids=nav_rids)
        rels_str = make_slide_rels(rid_map)
        files[f"ppt/slides/slide{snum}.xml"]           = xml_str.encode("utf-8")
        files[f"ppt/slides/_rels/slide{snum}.xml.rels"] = rels_str.encode("utf-8")

    # Mettre à jour presentation.xml (sldIdLst + relationships)
    prs_xml  = files["ppt/presentation.xml"].decode("utf-8")
    rels_xml = files["ppt/_rels/presentation.xml.rels"].decode("utf-8")

    for snum in range(FIRST_AGENCY_SLIDE, max_existing + 1):
        prs_xml  = re.sub(rf'\s*<p:sldId[^>]*r:id="rId{snum}"[^/]*/>', "", prs_xml)

    # Calculer dynamiquement les prochains rId et sldId libres
    # pour éviter toute collision avec les rIds existants du template
    import re as _re
    _existing_rids     = set(int(m) for m in _re.findall(r'Id="rId(\d+)"', rels_xml))
    _existing_sld_ids  = set(int(m) for m in _re.findall(r'<p:sldId id="(\d+)"', prs_xml))
    next_rid    = max(_existing_rids,    default=25) + 1
    next_sld_id = max(_existing_sld_ids, default=300) + 1

    # Mapping snum → rId (utilisé plus bas pour make_nav_rid_map)
    _snum_to_rid = {}

    new_sld_ids = ""
    new_rels    = ""
    for snum in sorted(slides.keys()):
        rid = f"rId{next_rid}"
        _snum_to_rid[snum] = rid
        new_sld_ids += f'\n    <p:sldId id="{next_sld_id}" r:id="{rid}"/>'
        new_rels    += (
            f'\n  <Relationship Id="{rid}" '
            f'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" '
            f'Target="slides/slide{snum}.xml"/>'
        )
        next_rid    += 1
        next_sld_id += 1


    prs_xml  = prs_xml.replace("</p:sldIdLst>", new_sld_ids + "\n  </p:sldIdLst>")
    rels_xml = rels_xml.replace("</Relationships>",  new_rels + "\n</Relationships>")

    files["ppt/presentation.xml"]            = prs_xml.encode("utf-8")
    files["ppt/_rels/presentation.xml.rels"] = rels_xml.encode("utf-8")

    # ── Mettre à jour [Content_Types].xml ─────────────────────────────────────
    # Les slides agency générées doivent y être déclarées sinon PowerPoint refuse
    # d'ouvrir le fichier ("cannot read")
    SLIDE_CT = 'application/vnd.openxmlformats-officedocument.presentationml.slide+xml'
    ct_xml   = files["[Content_Types].xml"].decode("utf-8")
    for snum in sorted(slides.keys()):
        part_name = f"/ppt/slides/slide{snum}.xml"
        if part_name not in ct_xml:
            entry = f'  <Override PartName="{part_name}" ContentType="{SLIDE_CT}"/>'
            ct_xml = ct_xml.replace("</Types>", entry + "\n</Types>")
    files["[Content_Types].xml"] = ct_xml.encode("utf-8")

    # Repack en PPTX
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zout:
        for name, data in files.items():
            zout.writestr(name, data)
    buf.seek(0)
    return buf.read()

# ─────────────────────────────────────────────────────────────
# 4. STANDALONE TEST
# ─────────────────────────────────────────────────────────────
if __name__ == "__main__":
    import sys

    EXCEL    = sys.argv[1] if len(sys.argv) > 1 else \
               "/mnt/user-data/uploads/Newbiz_Balance_DB_Report__1_.xlsx"
    TEMPLATE = sys.argv[2] if len(sys.argv) > 2 else \
               "T21_HK_Agencies_Glass_v12.pptx"
    OUTPUT   = sys.argv[3] if len(sys.argv) > 3 else \
               "/mnt/user-data/outputs/NBB_Agency_Cards_AUTO.pptx"

    print(f"📂 Excel    : {EXCEL}")
    print(f"📐 Template : {TEMPLATE}")

    df = pd.read_excel(EXCEL)
    print(f"✅ {len(df)} lignes chargées")

    result = build_agency_pptx(df, TEMPLATE)

    os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)
    with open(OUTPUT, "wb") as f:
        f.write(result)
    print(f"✅ Généré : {OUTPUT}  ({len(result)//1024} KB)")
