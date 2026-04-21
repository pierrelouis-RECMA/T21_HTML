import io, zipfile, pandas as pd
from pptx import Presentation
from pptx.util import Pt

def build_agency_pptx(df, template_path):
    # --- 1. ANALYSE DES DONNÉES ---
    df.columns = [c.strip() for c in df.columns]
    spend_col = "Integrated Spends" if "Integrated Spends" in df.columns else "Ad Spends"
    df[spend_col] = pd.to_numeric(df[spend_col], errors='coerce').fillna(0)

    agencies_stats = []
    for agency, group in df.groupby('Agency'):
        wins = group[group['NewBiz'] == 'WIN'][spend_col].sum()
        deps = group[group['NewBiz'] == 'DEPARTURE'][spend_col].sum()
        total_nbb = wins - abs(deps)
        
        # Filtrage > 3M pour les détails
        sig = group[group[spend_col].abs() >= 3]
        wins_txt = " · ".join([f"{r['Advertiser']} {r[spend_col]:.0f}m" for _, r in sig[sig['NewBiz'] == 'WIN'].iterrows()])
        deps_txt = " · ".join([f"{r['Advertiser']} {abs(r[spend_col]):.0f}m" for _, r in sig[sig['NewBiz'] == 'DEPARTURE'].iterrows()])

        agencies_stats.append({
            'agency': str(agency),
            'nbb': total_nbb,
            'wins_desc': wins_txt if wins_txt else "-",
            'deps_desc': deps_txt if deps_txt else "-"
        })

    # Tri par NBB
    nbb_ranking = sorted(agencies_stats, key=lambda x: x['nbb'], reverse=True)

    # --- 2. INJECTION DANS LE PPTX ---
    prs = Presentation(template_path)

    # Exemple : Mise à jour de la Slide 4 (Index 3)
    # On cherche les zones de texte pour y insérer nos résultats
    slide4 = prs.slides[3] 
    for shape in slide4.shapes:
        if shape.has_text_frame:
            # Ici on remplace des balises si elles existent ou on remplit des zones spécifiques
            # Pour l'exemple, on affiche les 3 premières agences
            if "Agency" in shape.text:
                top_text = ""
                for ag in nbb_ranking[:5]:
                    top_text += f"{ag['agency']}: {ag['nbb']:.1f}m (Wins: {ag['wins_desc']})\n"
                shape.text = top_text

    # --- 3. RETOUR DU FICHIER ---
    pptx_io = io.BytesIO()
    prs.save(pptx_io)
    return pptx_io.getvalue()